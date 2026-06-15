#!/usr/bin/env python3
"""Gera apresentação PPTX — Índices de Vegetação · A3 Computação Gráfica."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(BASE, "site", "apresentacao.pptx")

# ── Paleta ────────────────────────────────────────────────────────────────────
BG      = RGBColor(0xF0, 0xEE, 0xE6)   # warm cream
SURFACE = RGBColor(0xFA, 0xF9, 0xF5)   # card background
SURF2   = RGBColor(0xEA, 0xE7, 0xDC)   # table header
INK     = RGBColor(0x14, 0x14, 0x13)   # título
BODY    = RGBColor(0x45, 0x44, 0x3E)   # corpo
MUTED   = RGBColor(0x87, 0x86, 0x7F)   # legendas
CLAY    = RGBColor(0xD9, 0x77, 0x57)   # terracota accent
CLAY_D  = RGBColor(0xCC, 0x78, 0x5C)   # accent escuro (eyebrow)
BORDER  = RGBColor(0xE3, 0xDF, 0xD3)   # borda
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)


# ── Primitivos ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
_blank = prs.slide_layouts[6]


def slide():
    sl = prs.slides.add_slide(_blank)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return sl


def t(sl, text, l, top, w, h, sz=14, bold=False,
      color=INK, align=PP_ALIGN.LEFT, face="Calibri", italic=False):
    bx = sl.shapes.add_textbox(Inches(l), Inches(top), Inches(w), Inches(h))
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(sz)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    run.font.name     = face
    return bx


def eye(sl, text, l=0.75, top=0.38):
    t(sl, text.upper(), l, top, 11.8, 0.32, sz=9.5, bold=True, color=CLAY_D)


def h2(sl, text, l=0.75, top=0.75, w=11.8, h=1.1, sz=34):
    t(sl, text, l, top, w, h, sz=sz, color=INK, face="Georgia")


def img(sl, path, l, top, w, h):
    if not os.path.exists(path):
        print(f"  WARN missing: {path}")
        return
    sl.shapes.add_picture(path, Inches(l), Inches(top), Inches(w), Inches(h))


def outputs(name):
    return os.path.join(BASE, "outputs", name)


def images(name):
    return os.path.join(BASE, "images", name)


def box(sl, l, top, w, h, fill=SURFACE, border=BORDER, lw=0.5):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(top), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(lw)
    return sh


def card(sl, l, top, w, h, title="", body_text="", accent=False, num=""):
    bc = CLAY if accent else BORDER
    box(sl, l, top, w, h, border=bc)
    y = top + 0.17
    if num:
        t(sl, num, l+0.17, y, w-0.34, 0.6, sz=38, color=CLAY, face="Georgia")
        y += 0.62
    if title:
        t(sl, title, l+0.17, y, w-0.34, 0.38, sz=13, bold=True, color=INK)
        y += 0.4
    if body_text:
        t(sl, body_text, l+0.17, y, w-0.34, h-(y-top)-0.14, sz=12, color=BODY)


def pivot(sl, tag, text, l, top, w, h):
    box(sl, l, top, w, h, border=BORDER)
    bar = sl.shapes.add_shape(1, Inches(l), Inches(top), Inches(0.065), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = CLAY
    bar.line.color.rgb = CLAY; bar.line.width = Pt(0)
    if tag:
        t(sl, "↳ " + tag.upper(), l+0.15, top+0.1,  w-0.25, 0.27, sz=8.5, bold=True, color=CLAY_D)
    t(sl, text, l+0.15, top+(0.37 if tag else 0.12), w-0.25, h-0.5, sz=13, color=BODY)


def tbl(sl, headers, rows, l, top, col_w, row_h=0.46):
    for ri, row in enumerate([headers] + rows):
        is_h = ri == 0
        for ci, cell in enumerate(row):
            cx = l + sum(col_w[:ci])
            cy = top + ri * row_h
            box(sl, cx, cy, col_w[ci], row_h, fill=SURF2 if is_h else SURFACE)
            clr  = MUTED if is_h else BODY
            bold = is_h
            t(sl, str(cell), cx+0.1, cy+0.1, col_w[ci]-0.15, row_h-0.16,
              sz=10.5 if is_h else 12.5, bold=bold, color=clr)


# ═══════════════════════════════════════════════════════════════════════════════
# 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
bar = sl.shapes.add_shape(1, Inches(0.75), Inches(2.2), Inches(11.8), Inches(0.055))
bar.fill.solid(); bar.fill.fore_color.rgb = CLAY
bar.line.color.rgb = CLAY; bar.line.width = Pt(0)

t(sl, "AVALIAÇÃO A3 · COMPUTAÇÃO GRÁFICA · 2026.1 · NOTURNO",
  0.75, 1.78, 11.8, 0.36, sz=11, bold=True, color=CLAY_D)
t(sl, "Quanta vegetação existe nesta imagem?",
  0.75, 2.35, 11.8, 1.55, sz=50, color=INK, face="Georgia")
t(sl, "Um pipeline de processamento digital de imagens em Python — OpenCV, NumPy e "
      "Matplotlib — que transforma uma fotografia aérea RGB comum em uma resposta objetiva: "
      "quais pixels são vegetação e qual o percentual de cobertura. "
      "E que, no fim, se mostra maior do que o verde.",
  0.75, 4.05, 9.5, 1.4, sz=17, color=BODY)
t(sl, "Gilberto de Paiva Melo · Ciência da Computação",
  0.75, 5.7, 11.8, 0.5, sz=16, bold=True, color=INK)


# ═══════════════════════════════════════════════════════════════════════════════
# 2 — O PROBLEMA
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "O problema")
h2(sl, "O computador não vê verde. Vê números.")
t(sl, "O olho humano distingue vegetação sem esforço. Para a máquina, cada pixel é só um trio "
      "[R, G, B] de 0 a 255 — e é preciso um algoritmo para transformar esses números em "
      "informação ambiental.",
  0.75, 1.92, 11.8, 0.78, sz=16, color=BODY)

cw = 3.7
cy = 2.85
card(sl, 0.75,          cy, cw, 4.15, title="O objetivo",
     body_text="Medir o percentual de cobertura vegetal de uma cena automaticamente, "
               "separando planta de solo, estrada e céu.")
card(sl, 0.75+cw+0.2,   cy, cw, 4.15, title="Por que importa",
     body_text="Agricultura de precisão, monitoramento de desmatamento e áreas verdes urbanas — "
               "onde imagem de satélite multiespectral é cara e rara.")
card(sl, 0.75+2*(cw+0.2), cy, cw, 4.15, title="A restrição",
     body_text="Usar apenas câmera RGB comum — sem sensor infravermelho, sem equipamento "
               "especializado. Só Python e visão computacional.")


# ═══════════════════════════════════════════════════════════════════════════════
# 3 — AS IMAGENS
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "As imagens · mudança de caminho nº 1")
h2(sl, "Duas imagens, dois papéis")

img(sl, images("drone-campo-fazenda.jpg"), 0.75, 1.88, 5.75, 3.75)
t(sl, "drone-campo-fazenda · 19,9 MP · IMAGEM PRINCIPAL — campo + estrada + céu — todo o pipeline roda nela",
  0.75, 5.68, 5.75, 0.65, sz=12, color=MUTED)

img(sl, images("drone-lavoura-solo.jpg"), 6.7, 1.88, 5.75, 3.75)
t(sl, "drone-lavoura-solo · 11,9 MP · CASO DIFÍCIL — solo exposto que engana o índice — valida a segmentação",
  6.7, 5.68, 5.75, 0.65, sz=12, color=MUTED)

pivot(sl, "Mudança de caminho",
      "Começamos com foto de celular — mas o ângulo oblíquo misturava superfícies e degradava os "
      "índices. Migramos para imagens aéreas de drone em visada nadir, baixadas de bancos abertos "
      "com seleção documentada (seção 5.1 da proposta).",
      0.75, 6.38, 11.8, 0.87)


# ═══════════════════════════════════════════════════════════════════════════════
# 4 — PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Visão geral")
h2(sl, "O pipeline, em nove passos")

steps = [
    ("1", "Aquisição",             "cv2.imread() + conversão BGR→RGB e metadados"),
    ("2", "Pré-processamento",     "redimensionar 19,9 MP → 1600 px, escala de cinza, normalizar [0,1]"),
    ("3", "Histogramas + equaliz.", "distribuição R, G, B e realce de contraste global"),
    ("4", "Índices espectrais",    "VARI e ExG — a assinatura de cor da vegetação"),
    ("5", "Filtros",               "Gaussiano, Mediana e Bilateral, medidos por PSNR/SNR"),
    ("6", "Bordas (Canny)",        "3 pares de limiares comparados — 50/130 adotado"),
    ("7", "Segmentação",           "Otsu (brilho) vs. ExG (cor) + morfologia"),
    ("8", "Seg. aprimorada",       "ExGR + trava HSV — o verde bem identificado"),
    ("9", "Generalização",         "a mesma receita no petróleo (NASA) e na astronomia (Hubble)"),
]

cw, ch = 3.87, 1.52
for i, (n, title, sub) in enumerate(steps):
    col = i % 3
    row = i // 3
    l   = 0.75 + col * (cw + 0.17)
    top = 1.93 + row * (ch + 0.13)
    box(sl, l, top, cw, ch)
    circ = sl.shapes.add_shape(9, Inches(l+0.14), Inches(top+0.16), Inches(0.3), Inches(0.3))
    circ.fill.solid(); circ.fill.fore_color.rgb = CLAY
    circ.line.color.rgb = CLAY
    t(sl, n, l+0.14, top+0.14, 0.3, 0.33, sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t(sl, title, l+0.54, top+0.16, cw-0.68, 0.33, sz=12.5, bold=True, color=INK)
    t(sl, sub,   l+0.17, top+0.56, cw-0.3,  0.83, sz=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# 5 — METADADOS + CINZA
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passos 1–2 · descoberta nº 1")
h2(sl, "Brilho não identifica vegetação")

img(sl, outputs("original_vs_grayscale.png"), 0.75, 1.88, 8.0, 4.35)
t(sl, "Na versão cinza, campo, mata e solo se confundem — a cor era a informação",
  0.75, 6.27, 8.0, 0.48, sz=11.5, color=MUTED, align=PP_ALIGN.CENTER)

card(sl, 9.0, 1.88, 3.55, 2.0,
     title="Metadados da imagem",
     body_text="5464 × 3640 px (19,9 MP) · 3 canais · uint8 · 3,1 MB\n"
               "Valor médio: R 97 · G 118 · B 57")
pivot(sl, "Descoberta",
      "A conversão para cinza descarta a crominância — exatamente o que define a planta. "
      "Isso motivou os índices espectrais.",
      9.0, 4.05, 3.55, 1.7)


# ═══════════════════════════════════════════════════════════════════════════════
# 6 — HISTOGRAMAS
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 3")
h2(sl, "Lendo a cena pelos histogramas")

img(sl, outputs("histograma_rgb.png"), 0.75, 1.88, 5.75, 4.0)
t(sl, "Canais R, G, B — o verde desloca-se para intensidades mais altas",
  0.75, 5.93, 5.75, 0.48, sz=12, color=MUTED, align=PP_ALIGN.CENTER)

img(sl, outputs("histograma_exg.png"), 6.7, 1.88, 5.75, 4.0)
t(sl, "Índice ExG — massa concentrada acima de 0 (média 0,32): cena vegetada",
  6.7, 5.93, 5.75, 0.48, sz=12, color=MUTED, align=PP_ALIGN.CENTER)

t(sl, "O histograma do ExG já mostra o corte natural: pixels acima de zero têm verde em excesso — a base da segmentação.",
  0.75, 6.5, 11.8, 0.55, sz=15, color=BODY)


# ═══════════════════════════════════════════════════════════════════════════════
# 7 — ÍNDICES
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 4 · descoberta nº 2")
h2(sl, "VARI e ExG — a assinatura espectral do verde")

box(sl, 0.75, 1.83, 5.4, 0.72, fill=SURFACE)
t(sl, "VARI = (G − R) / (G + R − B)", 0.85, 1.9, 5.2, 0.55, sz=18, color=INK, face="Georgia", italic=True)

box(sl, 6.4, 1.83, 4.0, 0.72, fill=SURFACE)
t(sl, "ExG = 2G − R − B", 6.5, 1.9, 3.8, 0.55, sz=18, color=INK, face="Georgia", italic=True)

img(sl, outputs("comparativo_indices.png"), 0.75, 2.65, 11.8, 3.75)

pivot(sl, "Descoberta",
      "Na sombra, o denominador do VARI tende a zero e o índice explode até 144.707 "
      "(outliers sem significado físico). Solução: máscara de luminância — pixels com brilho < 0,1 são "
      "excluídos. A média mal muda (0,1392 → 0,1391), provando que era só ruído numérico.",
      0.75, 6.48, 11.8, 0.77)


# ═══════════════════════════════════════════════════════════════════════════════
# 8 — FILTROS
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 5")
h2(sl, "Três filtros, um vencedor")

img(sl, outputs("comparacao_filtros.png"), 0.75, 1.88, 11.8, 3.65)

headers = ["Filtro",            "PSNR (dB)", "SNR (dB)", "Leitura"]
rows    = [
    ["Gaussiano σ=1,5", "33,16", "26,06", "borra tudo por igual — menor fidelidade"],
    ["Mediana 5×5",     "32,81", "25,71", "altera textura fina do campo"],
    ["Bilateral d=9",   "34,08", "26,98", "suaviza preservando bordas — o mais fiel ✓"],
]
tbl(sl, headers, rows, l=0.75, top=5.65, col_w=[3.5, 1.55, 1.55, 5.15])


# ═══════════════════════════════════════════════════════════════════════════════
# 9 — CANNY
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 6")
h2(sl, "Canny: os limiares decidem o que é borda")

img(sl, outputs("canny_limiares.png"), 0.75, 1.88, 11.8, 3.35)

cw3, cy3 = 3.7, 5.33
card(sl, 0.75,            cy3, cw3, 1.95, title="30/90 — baixos",
     body_text="1,36% de bordas — captura até a textura da plantação. Ruidoso.")
card(sl, 0.75+cw3+0.17,   cy3, cw3, 1.95, title="50/130 — adotado ✓",
     body_text="0,71% — estrada, horizonte e talhões contínuos. Equilíbrio.", accent=True)
card(sl, 0.75+2*(cw3+0.17), cy3, cw3, 1.95, title="150/250 — altos",
     body_text="0,04% — só contornos fortes sobrevivem. Fragmentado.")


# ═══════════════════════════════════════════════════════════════════════════════
# 10 — OTSU VS EXG
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 7 · a grande descoberta")
h2(sl, "O Otsu marcou a estrada.\nO ExG marcou o campo.", sz=30)

img(sl, outputs("comparacao_segmentacao.png"), 0.75, 1.88, 7.8, 4.55)

card(sl, 8.78, 1.88, 3.78, 2.1, title="Otsu (brilho)",
     body_text="Limiar 0,55 segmenta a ESTRADA de terra clara — o que não é vegetação.")
t(sl, "17%", 9.05, 2.28, 2.0, 0.8, sz=40, color=CLAY, face="Georgia")

card(sl, 8.78, 4.1, 3.78, 2.1, title="ExG (cor) ✓",
     body_text="Segmenta pela assinatura espectral — campo verde e exclui estrada e céu.", accent=True)
t(sl, "83%", 9.05, 4.5, 2.0, 0.8, sz=40, color=CLAY, face="Georgia")

card(sl, 8.78, 6.35, 3.78, 0.9,
     body_text="Para vegetação RGB: crominância vence luminância. O Otsu pode segmentar o alvo errado.")


# ═══════════════════════════════════════════════════════════════════════════════
# 11 — VISUALIZAÇÃO
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Visualização · mudança de caminho nº 2")
h2(sl, "Como mostrar a detecção sem esconder a cena")

img(sl, outputs("segmentacao_exg.png"), 0.75, 1.88, 11.8, 4.35)
pivot(sl, "Mudança de caminho (em duas tentativas)",
      "1ª versão: pintar verde sólido — com 88% de cobertura virava um bloco estourado. "
      "2ª versão: verde translúcido — aceitável mas tingia a cena. "
      "Final: DESTAQUE POR DESSATURAÇÃO — a vegetação mantém a cor original e o resto vira cinza. "
      "O que está em cor é exatamente o que o algoritmo classificou.",
      0.75, 6.33, 11.8, 0.92)


# ═══════════════════════════════════════════════════════════════════════════════
# 12 — CASO DIFÍCIL
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 8 · o caso difícil")
h2(sl, "Quando os índices discordam, há um erro")

img(sl, outputs("analise_multi_imagem.png"), 0.75, 1.88, 11.8, 3.75)

cw2 = 5.7
card(sl, 0.75, 5.73, cw2, 1.53, title="campo-fazenda · índices concordam",
     body_text="ExG 87,9% e VARI médio +0,14 → estimativa confiável.")
card(sl, 0.75+cw2+0.17, 5.73, cw2, 1.53, title="lavoura-solo · índices divergem ⚠",
     body_text="ExG diz 97,7%, mas o VARI é NEGATIVO (−0,07). Solo marrom aparece em cor — incluído indevidamente.",
     accent=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 13 — RESOLUÇÃO
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 8 · a resolução")
h2(sl, "Histórico → final: o verde bem identificado")

img(sl, outputs("segmentacao_aprimorada.png"), 0.75, 1.88, 11.8, 3.75)

cw2 = 5.7
card(sl, 0.75, 5.73, cw2, 1.53, title="Como resolve",
     body_text="ExGR = 3G − 2,4R − B desconta o vermelho do solo; trava HSV rejeita céu e superfícies "
               "desbotadas. Fechamento antes da abertura preserva as fileiras finas da lavoura.")
card(sl, 0.75+cw2+0.17, 5.73, cw2, 1.53, title="Os números", accent=True,
     body_text="campo-fazenda: 87,9% → 63,8%   (céu e estrada fora)\n"
               "lavoura-solo:   97,7% → 40,2%   (solo fora, cultura verde preservada)")


# ═══════════════════════════════════════════════════════════════════════════════
# 14 — GRID COMPLETO
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Fecho da vegetação · todas as etapas num painel")
h2(sl, "O pipeline de vegetação, completo")

img(sl, outputs("grid_pipeline_completo.png"), 0.75, 1.88, 11.8, 4.92)
t(sl, "Original → cinza → 3 filtros → Canny → VARI → ExG → Otsu+morfologia → contornos ExG   ·   Próximo: e se o alvo não for verde?",
  0.75, 6.88, 11.8, 0.45, sz=11.5, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 15 — GENERALIZAÇÃO PETRÓLEO
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 9 · escalando a metodologia · mudança de caminho nº 3")
h2(sl, "A mesma receita, outro alvo:\npetróleo no Golfo do México", sz=30)

img(sl, outputs("generalizacao_petroleo.png"), 0.75, 1.88, 11.8, 4.28)
pivot(sl, "Melhorando a acurácia — de novo, cor vence brilho",
      "Imagem NASA/MODIS da Deepwater Horizon (24/05/2010, domínio público). "
      "1ª versão: só BRILHO — nuvens também são claras, entravam como falso positivo (7,1%). "
      "Assinatura real: COR QUENTE (tan) — R − B ≈ 20–25 para óleo vs ≈ 0 para água e nuvens. "
      "Com R−B ≥ 14 + mesmo refino morfológico: mancha COERENTE de 17,4%. A mesma lição da vegetação.",
      0.75, 6.26, 11.8, 1.0)


# ═══════════════════════════════════════════════════════════════════════════════
# 16 — APLICAÇÕES PETRÓLEO
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Setor de petróleo")
h2(sl, "Três aplicações diretas da metodologia")

cw3b, cy3b = 3.7, 1.93
card(sl, 0.75,             cy3b, cw3b, 3.45, title="Monitoramento de derrames",
     body_text="Contornos + área % da mancha em imagens de satélite ou sobrevoo — "
               "exatamente o experimento que acabamos de ver. Triagem rápida e barata antes de radar (SAR).")
card(sl, 0.75+cw3b+0.17,   cy3b, cw3b, 3.45, title="Faixa de dutos (right-of-way)",
     body_text="Os contornos de vegetação detectam mato invadindo a faixa de servidão de "
               "oleodutos — priorizando a manutenção preventiva quilômetro a quilômetro.")
card(sl, 0.75+2*(cw3b+0.17), cy3b, cw3b, 3.45, title="Vazamento por estresse vegetal",
     body_text="Solo contaminado estressa a vegetação antes de o óleo aflorar: queda de "
               "ExG/VARI ao longo da linha do duto = alerta precoce de vazamento.")

pivot(sl, "Limitação honesta",
      "Sem bandas térmicas/radar, água com sedimento pode confundir-se com óleo fino. "
      "A detecção RGB atua como primeira camada de triagem — aponta onde usar sensores especializados.",
      0.75, 5.53, 11.8, 1.17)


# ═══════════════════════════════════════════════════════════════════════════════
# 17 — GENERALIZAÇÃO ESTRELAS
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Passo 10 · a ideia original — espaço e telescópios")
h2(sl, "A mesma receita conta estrelas")

img(sl, outputs("generalizacao_estrelas.png"), 0.75, 1.88, 11.8, 4.28)
pivot(sl, "Vegetação mede cobertura · óleo mede área · astronomia conta objetos",
      "Imagem Hubble (NASA/ESA) do aglomerado M13 (domínio público). Assinatura = BRILHO. "
      "Limiar global perdia estrelas fracas (7.736). Versão final reusa 3 peças do projeto: "
      "CLAHE (equalização local) eleva as fracas; white top-hat (realce local) as extrai; "
      "e o mesmo Otsu escolhe o corte automático. O mesmo connectedComponentsWithStats CONTA: "
      "17.216 estrelas. No núcleo denso elas se fundem — contagem = limite inferior.",
      0.75, 6.26, 11.8, 1.0)


# ═══════════════════════════════════════════════════════════════════════════════
# 18 — SÍNTESE
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Síntese · a virada do projeto")
h2(sl, "Uma metodologia, três domínios")
t(sl, "Os mesmos cinco blocos — assinatura → limiar absoluto → morfologia → componentes conexos → medida — "
      "resolveram problemas de meio ambiente, indústria e ciência. Só muda a assinatura.",
  0.75, 1.87, 11.8, 0.65, sz=15, color=BODY)

theaders = ["",              "Vegetação",            "Petróleo",             "Astronomia"]
trows    = [
    ["Imagem",       "drone (RGB)",          "NASA/MODIS",           "Hubble (M13)"],
    ["Assinatura",   "verde (ExG/ExGR)",     "cor quente (R−B)",     "brilho (ponto claro)"],
    ["Limiar",       "absoluto / HSV",       "absoluto (R−B ≥ 14)",  "Otsu automático"],
    ["Mede",         "% cobertura",          "% área da mancha",     "contagem de objetos"],
    ["Resultado",    "63,8% / 40,2%",        "17,4%",                "17.216 estrelas"],
    ["Limite",       "cultura seca/esparsa", "sedimento vs. óleo",   "núcleo denso (fusão)"],
]
tbl(sl, theaders, trows, l=0.75, top=2.62, col_w=[2.45, 3.12, 3.12, 3.12], row_h=0.56)

t(sl, "Meta-lição: toda detecção por limiar tem um regime onde sinal e fundo se confundem — e isso é o mesmo nos três domínios.",
  0.75, 7.0, 11.8, 0.38, sz=12, italic=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# 19 — CONCLUSÕES
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
eye(sl, "Conclusões")
h2(sl, "Três lições do projeto")

cw3c, cy3c = 3.7, 1.93
card(sl, 0.75,              cy3c, cw3c, 3.35, num="1", title="Cor vence brilho",
     body_text="Vale na vegetação (ExG > Otsu, que pegou a estrada) e no óleo "
               "(cor quente R−B > brilho, que pegava nuvens). A cor é a assinatura; o brilho engana.")
card(sl, 0.75+cw3c+0.17,    cy3c, cw3c, 3.35, num="2", title="Índices divergentes denunciam erros",
     body_text="ExG alto + VARI negativo na lavoura = alarme de falso positivo. "
               "A segmentação final (ExGR + HSV + fechamento) resolveu: 97,7% → 40,2% de vegetação real.")
card(sl, 0.75+2*(cw3c+0.17), cy3c, cw3c, 3.35, num="3", title="A contribuição é a metodologia",
     body_text="Os mesmos blocos mediram % de vegetação, área da mancha de óleo (NASA) "
               "e contagem de 17.216 estrelas (Hubble). Vale para qualquer alvo.")

pivot(sl, "Entregáveis",
      "Notebook Jupyter com 15 etapas comentadas · Relatório técnico em LaTeX · "
      "Imagens geradas em outputs/ · Repositório Git com README de execução",
      0.75, 5.48, 11.8, 0.75)

t(sl, "Obrigado!", 0.75, 6.48, 11.8, 0.75, sz=30, color=INK, face="Georgia", italic=True)


# ── Salvar ───────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Salvo: {OUT}")
